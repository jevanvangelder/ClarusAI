from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import io
import base64
from PIL import Image

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF file"""
    try:
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        
        return '\n\n'.join(text)
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from Word document"""
    try:
        doc_file = io.BytesIO(file_bytes)
        doc = Document(doc_file)
        
        text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text.append(paragraph.text)
        
        return '\n\n'.join(text)
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from PowerPoint presentation"""
    try:
        ppt_file = io.BytesIO(file_bytes)
        prs = Presentation(ppt_file)
        
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        
        return '\n\n'.join(text)
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def extract_text_from_txt(file_bytes: bytes) -> str:
    """Extract text from plain text file"""
    try:
        return file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        try:
            return file_bytes.decode('latin-1')
        except Exception as e:
            return f"Error reading TXT: {str(e)}"

def parse_file(filename: str, file_bytes: bytes) -> dict:
    """
    Parse file based on extension
    Returns: {"text": str, "image": str | None, "type": str}
    """
    extension = filename.lower().split('.')[-1]
    
    # Images (NEW!)
    if extension in ['png', 'jpg', 'jpeg', 'gif', 'webp']:
        base64_image = image_to_base64(file_bytes, filename)
        return {
            "text": f"[Afbeelding: {filename}]",
            "image": base64_image,
            "type": "image"
        }
    
    # PDFs
    elif extension == 'pdf':
        text = extract_text_from_pdf(file_bytes)
        return {"text": text, "image": None, "type": "pdf"}
    
    # Word documents
    elif extension == 'docx':
        text = extract_text_from_docx(file_bytes)
        return {"text": text, "image": None, "type": "docx"}
    
    # PowerPoint
    elif extension == 'pptx':
        text = extract_text_from_pptx(file_bytes)
        return {"text": text, "image": None, "type": "pptx"}
    
    # Text files
    elif extension == 'txt':
        text = extract_text_from_txt(file_bytes)
        return {"text": text, "image": None, "type": "txt"}
    
    # Unsupported
    else:
        return {
            "text": f"Unsupported file type: {extension}",
            "image": None,
            "type": "unknown"
        }

def image_to_base64(file_bytes: bytes, filename: str) -> str:
    """Convert image to base64 string for GPT-4 Vision"""
    try:
        # Open image
        image = Image.open(io.BytesIO(file_bytes))
        
        # Resize if too large (max 2048x2048 for Vision)
        max_size = 2048
        if image.width > max_size or image.height > max_size:
            image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
        
        # Convert to RGB if needed
        if image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Save to bytes
        buffer = io.BytesIO()
        image.save(buffer, format='JPEG', quality=85)
        
        # Encode to base64
        img_bytes = buffer.getvalue()
        base64_str = base64.b64encode(img_bytes).decode('utf-8')
        
        return f"data:image/jpeg;base64,{base64_str}"
    
    except Exception as e:
        print(f"Error converting image to base64: {e}")
        return None